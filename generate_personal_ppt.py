# -*- coding: utf-8 -*-
"""
个人工作汇报 PPT 生成器（答辩版）
================================
生成：成员1-流量采集与流量特征提取模块汇报.pptx

设计原则：
  - 面向课堂答辩，而非源码阅读
  - 每页 4~6 条要点，每条不超过 20 字
  - 不出现函数名、类名、参数列表、变量名
  - 用流程图/架构图代替大段文字
  - 结构：做什么 → 怎么做 → 什么效果

需要 python-pptx: pip install python-pptx
"""

import os
import sys
from datetime import datetime

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# 颜色主题
# ============================================================================

C_PRIMARY   = RGBColor(0x1A, 0x56, 0xDB)   # 深蓝
C_GREEN     = RGBColor(0x2D, 0x7D, 0x46)   # 深绿
C_ORANGE    = RGBColor(0xE4, 0x6A, 0x1B)   # 橙色
C_DARK      = RGBColor(0x2D, 0x2D, 0x2D)   # 正文
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED     = RGBColor(0x88, 0x88, 0x88)
C_CARD_BG   = RGBColor(0xF5, 0xF7, 0xFD)
C_CARD_BDR  = RGBColor(0xDC, 0xE3, 0xF0)
C_DONE_BG   = RGBColor(0xEE, 0xF7, 0xF0)
C_FUTURE_BG = RGBColor(0xFD, 0xF4, 0xEC)

# ============================================================================
# 辅助函数
# ============================================================================

def _page_num(slide, n, total):
    tb = slide.shapes.add_textbox(Inches(8.5), Inches(7.1), Inches(1.2), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{n} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = C_MUTED
    p.alignment = PP_ALIGN.RIGHT


def _footer_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(10), Inches(0.03))
    s.fill.solid(); s.fill.fore_color.rgb = C_PRIMARY; s.line.fill.background()


def _title(slide, text, sub=None):
    """页面标题"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.05))
    s.fill.solid(); s.fill.fore_color.rgb = C_PRIMARY; s.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C_PRIMARY

    if sub:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.72), Inches(8.8), Inches(0.35))
        tb2.text_frame.paragraphs[0].text = sub
        tb2.text_frame.paragraphs[0].font.size = Pt(12)
        tb2.text_frame.paragraphs[0].font.color.rgb = C_MUTED


def _text(slide, x, y, w, h, lines, size=Pt(13), color=None, spacing=1.5):
    """多行文本"""
    if color is None: color = C_DARK
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = size; p.font.color.rgb = color
        p.space_after = Pt(6); p.line_spacing = Pt(size.pt * spacing)
    return tb


def _bullet(slide, x, y, icon, text, size=Pt(13)):
    """带图标要点"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(8.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{icon}  {text}"; p.font.size = size; p.font.color.rgb = C_DARK


def _placeholder(slide, x, y, w, h, label):
    """图片占位框"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFC)
    s.line.color.rgb = C_PRIMARY; s.line.width = Pt(1.5); s.line.dash_style = 2
    tf = s.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    p0.text = f"📷  {label}"; p0.font.size = Pt(10); p0.font.color.rgb = C_PRIMARY; p0.font.bold = True
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    p1.text = "（请在此处插入截图）"; p1.font.size = Pt(8); p1.font.color.rgb = C_MUTED


def _card(slide, x, y, w, h, icon, title, desc_lines):
    """圆角卡片"""
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = C_CARD_BG
    c.line.color.rgb = C_CARD_BDR; c.line.width = Pt(0.5)

    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.3), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{icon}  {title}"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = C_PRIMARY

    _text(slide, x + 0.15, y + 0.45, w - 0.3, h - 0.55, desc_lines, size=Pt(10), spacing=1.45)


def _flow_arrow(slide, x, y, w, color=None):
    """右箭头"""
    if color is None: color = C_PRIMARY
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.28))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def _flow_box(slide, x, y, w, h, text, color=None, text_color=None):
    """流程框"""
    if color is None: color = C_PRIMARY
    if text_color is None: text_color = C_WHITE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = text; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = text_color


def _down_arrow(slide, x, y):
    """下箭头(用 V 形)"""
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.2), Inches(0.35))
    s.fill.solid(); s.fill.fore_color.rgb = C_MUTED; s.line.fill.background()


def _section_label(slide, x, y, text, color=None):
    """分区标签"""
    if color is None: color = C_PRIMARY
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.6), Inches(0.32))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    p = s.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = text; p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = C_WHITE


# ============================================================================
# PPT 构建
# ============================================================================

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    TOTAL = 14
    n = [0]

    def nx(): n[0] += 1; return n[0]

    blank = prs.slide_layouts[6]

    # ========================================================================
    # SLIDE 1 — 封面
    # ========================================================================
    sl = prs.slides.add_slide(blank)

    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(3.2))
    bg.fill.solid(); bg.fill.fore_color.rgb = C_PRIMARY; bg.line.fill.background()

    tb = sl.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(8.4), Inches(1.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "AI-IDS 智能入侵检测系统"; p.font.size = Pt(38)
    p.font.bold = True; p.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph(); p2.text = "流量采集与流量特征提取模块"; p2.font.size = Pt(28)
    p2.font.color.rgb = C_WHITE; p2.space_before = Pt(8)

    tb2 = sl.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.5))
    tb2.text_frame.paragraphs[0].text = "课程设计个人工作汇报"
    tb2.text_frame.paragraphs[0].font.size = Pt(18)
    tb2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)

    _text(sl, 0.8, 3.6, 8.4, 1.5, [
        f"汇报日期：{datetime.now().strftime('%Y年%m月%d日')}",
        "项目名称：AI-IDS 智能入侵检测系统",
        "技术栈：Python · Scapy · XGBoost · MySQL",
    ], size=Pt(14))

    bt = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), Inches(10), Inches(0.35))
    bt.fill.solid(); bt.fill.fore_color.rgb = C_PRIMARY; bt.line.fill.background()
    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 2 — 目录
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "目录", "CONTENTS")
    _footer_bar(sl)

    items = [
        ("01", "项目背景与目标"),
        ("02", "系统整体架构"),
        ("03", "网络流量采集"),
        ("04", "协议深度解析"),
        ("05", "域名精准过滤"),
        ("06", "流量特征提取"),
        ("07", "持续监控引擎"),
        ("08", "数据库设计"),
        ("09", "AI 检测管线"),
        ("10", "端到端数据流"),
        ("11", "技术亮点总结"),
        ("12", "总结与展望"),
    ]
    for i, (no, title) in enumerate(items):
        y = 1.3 + i * 0.47
        c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y), Inches(0.35), Inches(0.35))
        c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY; c.line.fill.background()
        p = c.text_frame.paragraphs[0]; p.text = no; p.font.size = Pt(10)
        p.font.color.rgb = C_WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        _text(sl, 1.55, y + 0.02, 6, 0.3, [title], size=Pt(14))

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 3 — 项目背景与目标
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "项目背景与目标", "Background & Objectives")
    _footer_bar(sl)

    # 左侧：背景
    _section_label(sl, 0.6, 1.25, "项目背景")
    _text(sl, 0.6, 1.7, 4.2, 2.8, [
        "网络安全威胁日益复杂",
        "传统防火墙难以应对应用层攻击",
        "人工分析海量流量不现实",
        "需要自动化、智能化的检测方案",
        "课程设计要求完整的工程实践",
    ], size=Pt(12), spacing=1.6)

    # 右侧：目标
    _section_label(sl, 5.2, 1.25, "项目目标")
    _text(sl, 5.2, 1.7, 4.2, 2.8, [
        "构建完整的网络入侵检测系统",
        "实现实时流量采集与协议解析",
        "提取多维流量特征供 AI 分析",
        "结合统计模型与规则引擎精准判攻",
        "提供可视化管理平台便于运维",
    ], size=Pt(12), spacing=1.6)

    # 底部：我的分工
    _section_label(sl, 0.6, 4.7, "个人分工", C_GREEN)
    _text(sl, 0.6, 5.15, 8.8, 1.8, [
        "负责流量采集层的全部设计与开发",
        "涵盖网络抓包、协议解析、域名过滤、特征提取、持续监控五大模块",
        "设计 flow_features 数据库表结构并实现自动迁移机制",
        "开发 AI 持续检测管线，实现 XGBoost + 正则双引擎融合决策",
        "总计 7 个核心模块，5000+ 行代码",
    ], size=Pt(11), spacing=1.5)

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 4 — 系统整体架构
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "系统整体架构", "System Architecture")
    _footer_bar(sl)

    # 四层架构图
    layers = [
        ("数据采集层", C_PRIMARY,
         ["实时网络抓包", "协议深度解析", "域名精准过滤", "流特征提取", "持续监控编排"]),
        ("AI 推理层", RGBColor(0x6C, 0x34, 0x80),
         ["统计模型推理", "正则载荷扫描", "双引擎融合决策", "自动生成判定依据"]),
        ("Web 展示层", C_GREEN,
         ["实时仪表盘", "攻击告警列表", "流量趋势图表", "IP 封禁管理"]),
        ("基础设施层", C_ORANGE,
         ["MySQL 数据存储", "Npcap 抓包驱动", "Flask Web 框架", "ECharts 可视化"]),
    ]

    for i, (name, color, items) in enumerate(layers):
        y = 1.25 + i * 1.4
        # 层标签
        lb = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(1.6), Inches(1.1))
        lb.fill.solid(); lb.fill.fore_color.rgb = color; lb.line.fill.background()
        tf = lb.text_frame; tf.word_wrap = True
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        p0.text = name; p0.font.size = Pt(13); p0.font.bold = True; p0.font.color.rgb = C_WHITE
        p0.space_before = Pt(12)

        # 内容卡片
        for j, item in enumerate(items):
            bx = 2.4 + j * 1.5
            _flow_box(sl, bx, y + 0.1, 1.35, 0.42, item, color=C_CARD_BG, text_color=C_DARK)

        # 层间箭头
        if i < 3:
            _flow_arrow(sl, 8.8, y + 0.45, 0.5, C_MUTED)

    # 底部说明
    _text(sl, 0.5, 6.95, 9, 0.3,
          ["四层架构各司其职，层间通过标准数据结构传递，支持独立测试与并行开发"],
          size=Pt(9), color=C_MUTED)

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 5 — 网络流量采集
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "网络流量采集", "Traffic Capture")
    _footer_bar(sl)

    # 左侧 — 功能要点
    _section_label(sl, 0.5, 1.25, "核心能力")
    _bullet(sl, 0.5, 1.75, "📡", "实时监听网卡流量")
    _bullet(sl, 0.5, 2.25, "🔍", "支持 TCP / UDP / ICMP 协议")
    _bullet(sl, 0.5, 2.75, "🎯", "自定义过滤规则与目标监控")
    _bullet(sl, 0.5, 3.25, "💾", "数据自动缓存与 CSV 导出")
    _bullet(sl, 0.5, 3.75, "🔒", "线程安全设计，后台持续运行")

    _section_label(sl, 0.5, 4.45, "解决什么问题")
    _text(sl, 0.5, 4.9, 4.5, 1.8, [
        "为整个 IDS 系统提供原始数据源",
        "在系统底层完成流量筛选，减少后续处理压力",
        "输出标准化数据格式供下游模块消费",
    ], size=Pt(11), spacing=1.5)

    # 右侧 — 流程示意图
    _section_label(sl, 5.3, 1.25, "工作流程")
    _flow_box(sl, 6.0, 1.85, 3.0, 0.55, "网卡原始流量")
    _down_arrow(sl, 7.4, 2.42)
    _flow_box(sl, 6.0, 2.85, 3.0, 0.55, "BPF 内核过滤")
    _down_arrow(sl, 7.4, 3.42)
    _flow_box(sl, 6.0, 3.85, 3.0, 0.55, "数据包结构化解析")
    _down_arrow(sl, 7.4, 4.42)
    _flow_box(sl, 6.0, 4.85, 3.0, 0.55, "缓存 + CSV 导出")

    # 截图占位
    _placeholder(sl, 5.3, 5.8, 4.2, 1.2, "流量采集运行截图\n（终端输出 + 统计信息）")

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 6 — 协议深度解析
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "协议深度解析", "Protocol Parsing")
    _footer_bar(sl)

    _section_label(sl, 0.5, 1.25, "支持协议")
    _text(sl, 0.5, 1.65, 4.5, 1.5, [
        "网络层：IP（含 TTL 提取）",
        "传输层：TCP · UDP · ICMP",
        "应用层：HTTP（请求方法 + 头部字段）",
        "加密层：TLS（握手类型 + SNI 提取）",
    ], size=Pt(12), spacing=1.6)

    _section_label(sl, 0.5, 3.35, "解析能力")
    _text(sl, 0.5, 3.75, 4.5, 2.2, [
        "HTTP：提取请求方法、Host、URI、User-Agent 等完整头部",
        "TLS：识别 ClientHello / ServerHello，",
        "     按 RFC 5246 标准逐字节提取 SNI 主机名",
        "TCP：解析 SYN/ACK/FIN/RST/PSH/URG 六种标志位",
        "ICMP：区分 Echo Request 与 Echo Reply",
    ], size=Pt(11), spacing=1.45)

    _section_label(sl, 5.3, 1.25, "解析流程")
    _flow_box(sl, 6.0, 1.75, 3.0, 0.45, "原始数据包")
    _down_arrow(sl, 7.4, 2.22)
    _flow_box(sl, 6.0, 2.6, 3.0, 0.45, "IP 层提取")
    _down_arrow(sl, 7.4, 3.07)
    _flow_box(sl, 5.0, 3.45, 1.45, 0.45, "TCP 解析")
    _flow_box(sl, 6.55, 3.45, 1.45, 0.45, "UDP 解析")
    _flow_box(sl, 8.1, 3.45, 1.45, 0.45, "ICMP 解析")
    _down_arrow(sl, 7.4, 3.92)
    _flow_box(sl, 5.5, 4.3, 2.0, 0.45, "HTTP 检测")
    _flow_box(sl, 7.7, 4.3, 2.0, 0.45, "TLS 检测")
    _down_arrow(sl, 7.4, 4.77)
    _flow_box(sl, 6.0, 5.15, 3.0, 0.45, "输出统一解析结果")

    _placeholder(sl, 5.3, 5.9, 4.2, 1.1, "协议解析运行截图\n（终端表格输出示例）")

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 7 — 域名精准过滤
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "域名精准过滤", "Domain-Based Filtering")
    _footer_bar(sl)

    _section_label(sl, 0.5, 1.25, "为什么需要域名过滤")
    _text(sl, 0.5, 1.65, 5.5, 1.2, [
        "目标网站 IP 可能动态变化（CDN / 负载均衡）",
        "仅靠 IP 过滤会丢失 TLS 加密流量中的域名信息",
        "本地部署场景需要双向流量完整捕获",
    ], size=Pt(11), spacing=1.5)

    _section_label(sl, 0.5, 2.95, "三级过滤策略")
    _text(sl, 0.5, 3.35, 5.5, 2.8, [
        "第一级 — 内核层预过滤",
        "  域名 DNS 解析为 IP → 注入内核 BPF 过滤器",
        "  在内核层面尽早丢弃无关流量",
        "",
        "第二级 — HTTP Host 头匹配",
        "  扫描数据载荷中的 Host 字段",
        "  支持子域名后缀匹配（如 api.example.com）",
        "",
        "第三级 — TLS SNI 扩展匹配",
        "  解析 TLS ClientHello 握手包",
        "  提取 SNI 扩展中的目标主机名进行比对",
    ], size=Pt(11), spacing=1.3)

    _section_label(sl, 3.2, 6.25, "效果")
    _text(sl, 3.2, 6.5, 4, 0.5,
          ["过滤后仅保留目标网站相关流量，大幅降低后续处理压力"],
          size=Pt(11), color=C_GREEN)

    # 右侧流程图
    _section_label(sl, 6.2, 1.25, "过滤流程")
    _flow_box(sl, 6.5, 1.8, 2.8, 0.4, "收到数据包")
    _down_arrow(sl, 7.8, 2.22)
    _flow_box(sl, 6.5, 2.55, 2.8, 0.4, "Tier1: BPF IP 过滤")
    _down_arrow(sl, 7.8, 2.97)
    _flow_box(sl, 6.5, 3.3, 2.8, 0.4, "Tier2: HTTP Host 匹配")
    _down_arrow(sl, 7.8, 3.72)
    _flow_box(sl, 6.5, 4.05, 2.8, 0.4, "Tier3: TLS SNI 匹配")
    _down_arrow(sl, 7.8, 4.47)
    _flow_box(sl, 6.5, 4.8, 2.8, 0.4, "连接追踪缓存")
    _down_arrow(sl, 7.8, 5.22)
    _flow_box(sl, 6.5, 5.55, 2.8, 0.4, "决策：放行 / 丢弃", C_GREEN)

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 8 — 流量特征提取
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "流量特征提取", "Feature Extraction")
    _footer_bar(sl)

    _section_label(sl, 0.5, 1.25, "做什么")
    _text(sl, 0.5, 1.65, 5.0, 1.2, [
        "将逐包数据聚合为双向流（Flow）记录",
        "从原始数据包中提炼出 AI 可理解的特征向量",
        "一个流 = 一次完整的网络会话",
    ], size=Pt(12), spacing=1.55)

    _section_label(sl, 0.5, 3.0, "怎么做 — 双向流聚合")
    _text(sl, 0.5, 3.35, 5.0, 1.5, [
        "按五元组（源IP + 目的IP + 源端口 + 目的端口 + 协议）建立会话",
        "自动归一化双向流量：A→B 与 B→A 归入同一条流",
        "实时统计：包数、字节数、到达间隔、TCP 标志等",
        "流超时自动关闭（60 秒无活动）或 FIN/RST 触发结束",
    ], size=Pt(11), spacing=1.4)

    _section_label(sl, 0.5, 4.95, "输出 — 15 维特征向量")
    # 特征卡片
    feat_groups = [
        ("基础统计 (6维)", "协议号 · 持续时间\n前/后向包数\n前/后向最大包长"),
        ("流量速率 (4维)", "前/后向平均包长\n字节速率 · 包速率"),
        ("时序特征 (3维)", "前向到达间隔均值\n后向到达间隔均值"),
        ("标志统计 (2维)", "SYN 标志计数\nFIN/RST 标志计数"),
    ]
    for i, (name, desc) in enumerate(feat_groups):
        x = 0.5 + i * 2.3
        _card(sl, x, 5.35, 2.15, 1.35, "", name, desc.split("\n"))

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 9 — 持续监控引擎
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "持续监控引擎", "Continuous Monitoring Engine")
    _footer_bar(sl)

    _section_label(sl, 0.5, 1.25, "核心职责")
    _text(sl, 0.5, 1.65, 5.5, 1.8, [
        "串联采集 → 解析 → 特征提取 → 数据库全流程",
        "实现无人值守的长期后台运行",
        "自动处理 DNS 变更、数据库断连等异常",
    ], size=Pt(12), spacing=1.55)

    _section_label(sl, 0.5, 3.55, "运行机制")
    _text(sl, 0.5, 3.9, 5.5, 2.5, [
        "启动阶段：DNS 解析目标 → 初始化数据库 → 启动抓包",
        "运行循环：每秒检查新数据包 → 增量解析 → 送入流管理器",
        "定时写入：每 60 秒回收完成的流 → 批量写入 MySQL",
        "异常恢复：数据库断连自动重试，DNS 记录变更自动刷新",
        "优雅退出：停止抓包 → 最终数据回收 → 打印运行摘要",
    ], size=Pt(11), spacing=1.4)

    # 右侧截图
    _placeholder(sl, 6.3, 1.25, 3.3, 2.3,
                 "监控引擎运行截图\n（控制台实时状态面板）")
    _placeholder(sl, 6.3, 3.85, 3.3, 2.3,
                 "数据库表数据截图\n（flow_features 查询结果）")

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 10 — 数据库设计
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "数据库设计", "Database — flow_features 表")
    _footer_bar(sl)

    _section_label(sl, 0.5, 1.25, "设计思路")
    _text(sl, 0.5, 1.65, 5.0, 1.5, [
        "一张表承载完整流量画像 — 22 个字段覆盖",
        "元信息 · 统计特征 · 载荷数据三类字段分层",
        "4 个索引加速时间/主机/IP/协议维度查询",
    ], size=Pt(12), spacing=1.55)

    # 三列卡片
    col_data = [
        ("📋 元信息 (3列)", ["自增主键 + 创建时间", "目标主机域名", "目标 IP 地址"]),
        ("📊 统计特征 (15列)", ["基础流量统计 (6维)", "速率与均值 (4维)", "时序与标志 (5维)"]),
        ("📦 载荷数据 (5列)", ["HTTP 原始载荷", "请求 URI 与方法", "Host + User-Agent"]),
    ]
    for i, (title, items) in enumerate(col_data):
        _card(sl, 0.5 + i * 3.1, 3.3, 2.9, 1.7, "", title, items)

    _section_label(sl, 0.5, 5.25, "工程特性")
    _text(sl, 0.5, 5.6, 5.0, 1.2, [
        "自动迁移：检测已有列 → 仅补充缺失列，升级零停机",
        "AI 状态追踪：ai_processed + predict_time 字段标记分析进度",
        "支持正则引擎：保留 HTTP 载荷供特征库精准匹配",
    ], size=Pt(11), spacing=1.4)

    _placeholder(sl, 6.0, 5.25, 3.5, 1.6,
                 "数据库 ER 图\n（flow_features 表结构）")

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 11 — AI 检测管线
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "AI 检测管线", "AI Detection Pipeline")
    _footer_bar(sl)

    # 流程图
    _section_label(sl, 0.5, 1.25, "检测流程")

    steps = [
        ("数据库轮询", "每 60 秒查询\n未分析的新记录"),
        ("统计模型推理", "XGBoost\n输出攻击类型\n与置信度"),
        ("正则载荷扫描", "特征库匹配\nHTTP 载荷中的\n攻击签名"),
        ("融合决策", "双引擎结果\n智能仲裁\n输出最终判定"),
        ("结果写入", "存入告警日志\n标记已分析\nWeb 可实时查看"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = 0.4 + i * 1.9
        _flow_box(sl, x, 1.8, 1.7, 0.45, title)
        _text(sl, x + 0.05, 2.35, 1.6, 1.0, desc.split("\n"), size=Pt(8.5), spacing=1.3)
        if i < 4:
            _flow_arrow(sl, x + 1.75, 1.9, 0.1)

    _section_label(sl, 0.5, 3.5, "融合策略")
    _text(sl, 0.5, 3.9, 5.5, 1.5, [
        "仅统计模型命中 → 输出模型判定的攻击类型",
        "仅正则命中 → 以正则结果覆盖（更精确的攻击分类）",
        "两者同时命中 → 正则类型优先，双来源标注",
        "自动生成中文判定理由，解释检测依据",
    ], size=Pt(12), spacing=1.5)

    _section_label(sl, 0.5, 5.55, "运行模式")
    _text(sl, 0.5, 5.95, 5.5, 1.0, [
        "持续模式：后台常驻运行，适用于生产环境",
        "单次模式：检测一次后退出，适用于调试与演示",
        "全量重检：重置所有记录后重新分析，适用于模型升级",
    ], size=Pt(11), spacing=1.45)

    _placeholder(sl, 6.3, 3.5, 3.3, 2.0,
                 "AI 检测运行截图\n（终端输出: 预测类型 / 置信度 / 风险等级）")
    _placeholder(sl, 6.3, 5.75, 3.3, 1.2,
                 "告警日志表截图\n（traffic_logs 查询结果）")

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 12 — 端到端数据流
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "端到端数据流", "End-to-End Data Flow")
    _footer_bar(sl)

    # 横向流程图（两行）
    row1 = [
        ("网卡流量", C_PRIMARY),
        ("实时抓包", C_PRIMARY),
        ("协议解析", C_PRIMARY),
        ("域名过滤", C_PRIMARY),
        ("特征提取", C_PRIMARY),
    ]
    row2 = [
        ("监控引擎\n写入数据库", C_GREEN),
        ("AI 统计模型\n推理", RGBColor(0x6C, 0x34, 0x80)),
        ("正则引擎\n载荷扫描", RGBColor(0x6C, 0x34, 0x80)),
        ("融合决策\n写入告警", C_ORANGE),
        ("Web 仪表盘\n实时展示", C_ORANGE),
    ]

    for i, (text, color) in enumerate(row1):
        x = 0.4 + i * 1.9
        _flow_box(sl, x, 1.3, 1.55, 0.65, text, color)
        if i < 4:
            _flow_arrow(sl, x + 1.6, 1.48, 0.25)

    # 大向下箭头
    for i in range(5):
        _down_arrow(sl, 1.05 + i * 1.9, 2.0)

    for i, (text, color) in enumerate(row2):
        x = 0.4 + i * 1.9
        _flow_box(sl, x, 2.55, 1.55, 0.65, text, color)
        if i < 4:
            _flow_arrow(sl, x + 1.6, 2.73, 0.25)

    # 三阶段说明
    _section_label(sl, 0.5, 3.7, "阶段一：数据采集")
    _text(sl, 0.5, 4.05, 3.0, 1.2, [
        "网卡 → 抓包 → 解析",
        "→ 过滤 → 特征提取",
        "→ 写入数据库",
        "全流程约 60 秒/批次",
    ], size=Pt(10), spacing=1.45)

    _section_label(sl, 3.6, 3.7, "阶段二：AI 分析")
    _text(sl, 3.6, 4.05, 3.0, 1.2, [
        "轮询未分析记录",
        "双引擎并行推理",
        "智能融合决策",
        "写入告警日志",
    ], size=Pt(10), spacing=1.45)

    _section_label(sl, 6.7, 3.7, "阶段三：展示响应")
    _text(sl, 6.7, 4.05, 2.8, 1.2, [
        "Web 仪表盘实时刷新",
        "攻击趋势可视化",
        "支持 IP 一键封禁",
        "日志检索与回溯",
    ], size=Pt(10), spacing=1.45)

    # 底部关键数字
    stats = [("7", "核心模块"), ("5000+", "行代码"), ("15", "维特征"), ("22", "数据库字段"), ("60s", "检测周期")]
    for i, (num, label) in enumerate(stats):
        x = 0.5 + i * 1.9
        s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.5), Inches(1.6), Inches(0.85))
        s.fill.solid(); s.fill.fore_color.rgb = C_CARD_BG; s.line.color.rgb = C_PRIMARY; s.line.width = Pt(0.5)
        tf = s.text_frame; tf.word_wrap = True
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        p0.text = num; p0.font.size = Pt(24); p0.font.bold = True; p0.font.color.rgb = C_PRIMARY
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        p1.text = label; p1.font.size = Pt(9); p1.font.color.rgb = C_DARK

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 13 — 技术亮点总结
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "技术亮点总结", "Technical Highlights")
    _footer_bar(sl)

    highlights = [
        ("🧩", "模块化分层架构",
         "采集→解析→特征→AI→展示，层层解耦。\n每层可独立开发、测试、替换。"),
        ("📐", "双向流归一化",
         "自动将 A→B 与 B→A 映射为同一会话。\n消除方向歧义，AI 分析更准确。"),
        ("🔀", "双引擎融合决策",
         "统计模型负责异常发现，\n正则引擎负责精准识别。\n两者互补，降低误报与漏报。"),
        ("🔄", "增量统计算法",
         "无需缓存全部数据包即可计算\n均值、速率、到达间隔等特征。\n内存占用恒定，适合长期运行。"),
        ("📋", "数据库自动迁移",
         "表结构升级时自动检测并补充缺失列。\n无需手动执行 SQL，部署零停机。"),
        ("🛡️", "全链路异常容错",
         "从网络抓包到数据库写入，\n每个环节均有自动恢复机制。\n单点故障不影响整体运行。"),
        ("⏱️", "域名三级过滤",
         "内核 BPF → HTTP Host → TLS SNI\n逐级精准过滤，大幅减少无效处理。"),
        ("🧪", "完善的自测体系",
         "每个模块内置可运行的自测代码。\n无需外部依赖即可验证功能正确性。"),
    ]

    for i, (icon, title, desc) in enumerate(highlights):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 4.8
        y = 1.25 + row * 1.45

        # 图标 + 标题
        _text(sl, x, y, 4.5, 0.25, [f"{icon}  {title}"], size=Pt(12), color=C_PRIMARY)

        # 描述
        _text(sl, x + 0.25, y + 0.35, 4.2, 1.0, desc.split("\n"), size=Pt(9.5), spacing=1.3)

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # SLIDE 14 — 总结与展望
    # ========================================================================
    sl = prs.slides.add_slide(blank)
    _title(sl, "总结与展望", "Summary & Future Work")
    _footer_bar(sl)

    # 左侧已完成
    left = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.2), Inches(4.4), Inches(4.8))
    left.fill.solid(); left.fill.fore_color.rgb = C_DONE_BG
    left.line.color.rgb = C_GREEN; left.line.width = Pt(1)

    _text(sl, 0.6, 1.35, 4.0, 0.3, ["✅  已完成工作"], size=Pt(16), color=C_GREEN)
    _text(sl, 0.6, 1.8, 4.0, 4.0, [
        "📡 流量采集层",
        "   · 实时网络抓包引擎",
        "   · 多协议深度解析（HTTP/TLS/ICMP）",
        "   · 三级域名精准过滤",
        "   · 15 维流特征提取",
        "   · 全流程持续监控引擎",
        "   · 测试流量生成工具",
        "",
        "🤖 AI 推理层",
        "   · XGBoost + 正则双引擎检测管线",
        "   · 智能融合决策与自动判定理由生成",
        "",
        "🗄️ 数据层",
        "   · flow_features 表设计与自动迁移",
    ], size=Pt(10), spacing=1.35)

    # 右侧展望
    right = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2), Inches(4.4), Inches(4.8))
    right.fill.solid(); right.fill.fore_color.rgb = C_FUTURE_BG
    right.line.color.rgb = C_ORANGE; right.line.width = Pt(1)

    _text(sl, 5.4, 1.35, 4.0, 0.3, ["🔮  后续计划"], size=Pt(16), color=C_ORANGE)
    _text(sl, 5.4, 1.8, 4.0, 4.0, [
        "📡 采集层增强",
        "   · 支持 IPv6 流量采集",
        "   · 增加 DNS/DHCP 协议解析",
        "",
        "📊 特征层优化",
        "   · 引入熵、方差等高级统计特征",
        "   · 增加多时间窗口聚合",
        "",
        "🤖 AI 引擎升级",
        "   · 集成在线增量学习",
        "   · 扩充正则特征规则库",
        "   · 引入隔离森林等异常检测模型",
        "",
        "🏗️ 工程化完善",
        "   · Docker 容器化部署",
        "   · 性能压测与优化",
    ], size=Pt(10), spacing=1.35)

    # 感谢语
    tb = sl.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = "🙏  感谢各位老师和同学的指导与帮助！"
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = C_PRIMARY

    _page_num(sl, nx(), TOTAL)

    # ========================================================================
    # 保存
    # ========================================================================
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "成员1-流量采集与流量特征提取模块汇报.pptx")
    prs.save(out)
    print(f"\n✅ PPT 已生成: {out}")
    print(f"   共 {TOTAL} 页")
    return out


# ============================================================================
# 截图指南
# ============================================================================

GUIDE = """
╔══════════════════════════════════════════════════════════════════════════╗
║                      📷  需要截取/绘制的图片清单                          ║
╚══════════════════════════════════════════════════════════════════════════╝

共 6 处占位框需要替换：

┌──────────────────────────────────────────────────────────────────────────┐
│ Slide 5  — 流量采集运行截图                                              │
│ 运行: cd src\\traffic_module && python capture.py -t 30                   │
│ 注意: 需管理员权限 + Npcap                                                │
└──────────────────────────────────────────────────────────────────────────┘

┌──────────────────────────────────────────────────────────────────────────┐
│ Slide 6  — 协议解析运行截图                                              │
│ 运行: cd src\\traffic_module && python parser.py -n 50                    │
│ 无需抓包也可运行（自动使用模拟数据）                                      │
└──────────────────────────────────────────────────────────────────────────┘

┌──────────────────────────────────────────────────────────────────────────┐
│ Slide 9  — 两处截图                                                       │
│ A. 监控引擎运行: cd src\\traffic_module                                    │
│    python traffic_monitor.py --target www.baidu.com --interval 30         │
│ B. 数据库查询: SELECT * FROM flow_features ORDER BY id DESC LIMIT 20;     │
└──────────────────────────────────────────────────────────────────────────┘

┌──────────────────────────────────────────────────────────────────────────┐
│ Slide 10 — 数据库 ER 图                                                  │
│ 使用 DBeaver → 右键表 → View Diagram 截图                                 │
└──────────────────────────────────────────────────────────────────────────┘

┌──────────────────────────────────────────────────────────────────────────┐
│ Slide 11 — 两处截图                                                       │
│ A. AI 检测运行: cd src\\ai_engine && python run_pipeline.py --once        │
│ B. 告警日志: SELECT * FROM traffic_logs ORDER BY id DESC LIMIT 20;        │
└──────────────────────────────────────────────────────────────────────────┘

🛠️  截图准备步骤：

1. 确保 Python 环境就绪: pip install scapy pymysql
2. Windows 安装 Npcap: https://npcap.com/
3. 启动 MySQL 服务
4. 初始化数据库: cd web\\backend && python init_db.py
5. 完整演示 (3个终端):
   终端1: python src\\traffic_module\\traffic_monitor.py --target www.baidu.com
   终端2: python src\\traffic_module\\generate_test_traffic.py --target www.baidu.com
   终端3: python src\\ai_engine\\run_pipeline.py --once
6. 数据库工具推荐: DBeaver (免费) https://dbeaver.io/
"""

if __name__ == "__main__":
    print(GUIDE)
    try:
        from pptx import Presentation as _  # noqa: F401
    except ImportError:
        print("❌ 请先安装 python-pptx: pip install python-pptx")
        sys.exit(1)
    build()
